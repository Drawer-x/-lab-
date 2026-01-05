"""
自动生成演示PPT - 电影评论情感分析系统
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import os

# 颜色定义
DARK_BLUE = RGBColor(0x1a, 0x23, 0x5a)
WHITE = RGBColor(0xff, 0xff, 0xff)
GOLD = RGBColor(0xff, 0xd7, 0x00)
LIGHT_BLUE = RGBColor(0x4a, 0x6c, 0xf7)
GREEN = RGBColor(0x2e, 0xcc, 0x71)
RED = RGBColor(0xe7, 0x4c, 0x3c)


def add_background(slide):
    """添加深蓝色背景"""
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.33), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = DARK_BLUE
    bg.line.fill.background()
    spTree = slide.shapes._spTree
    sp = bg._element
    spTree.remove(sp)
    spTree.insert(2, sp)


def add_title(slide, text, top=Inches(0.8), size=48):
    """添加标题"""
    box = slide.shapes.add_textbox(Inches(0.5), top, Inches(12.33), Inches(1.2))
    p = box.text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER


def add_subtitle(slide, text, top=Inches(2)):
    """添加副标题"""
    box = slide.shapes.add_textbox(Inches(0.5), top, Inches(12.33), Inches(0.8))
    p = box.text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(28)
    p.font.color.rgb = GOLD
    p.alignment = PP_ALIGN.CENTER


def add_content(slide, items, top=Inches(2.2), size=28):
    """添加内容列表"""
    box = slide.shapes.add_textbox(Inches(1.2), top, Inches(11), Inches(5))
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = "• " + item
        p.font.size = Pt(size)
        p.font.color.rgb = WHITE
        p.space_after = Pt(20)


def add_table(slide, data, left=Inches(1.5), top=Inches(2.2)):
    """添加表格"""
    rows, cols = len(data), len(data[0])
    width = Inches(10)
    height = Inches(0.6 * rows)
    
    table = slide.shapes.add_table(rows, cols, left, top, width, height).table
    
    col_width = width / cols
    for i in range(cols):
        table.columns[i].width = int(col_width)
    
    for i, row in enumerate(data):
        for j, cell_text in enumerate(row):
            cell = table.cell(i, j)
            cell.text = str(cell_text)
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(20)
            p.alignment = PP_ALIGN.CENTER
            
            if i == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = LIGHT_BLUE
                p.font.color.rgb = WHITE
                p.font.bold = True
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(0x2a, 0x2a, 0x4a)
                p.font.color.rgb = WHITE


def add_code(slide, code, top=Inches(2.2)):
    """添加代码块"""
    # 背景框
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, 
                                    Inches(1), top, Inches(11.33), Inches(4.5))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0x1e, 0x1e, 0x1e)
    shape.line.color.rgb = RGBColor(0x3a, 0x3a, 0x3a)
    
    # 代码文本
    box = slide.shapes.add_textbox(Inches(1.3), top + Inches(0.3), Inches(10.7), Inches(4))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = code
    p.font.size = Pt(18)
    p.font.color.rgb = GREEN
    p.font.name = "Consolas"


def create_ppt():
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    
    # ===== 封面 =====
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_title(slide, "基于深度学习的", top=Inches(2), size=40)
    add_title(slide, "电影评论情感分析系统", top=Inches(2.8), size=48)
    add_subtitle(slide, "数据科学期末大作业", top=Inches(4.2))
    
    box = slide.shapes.add_textbox(Inches(0.5), Inches(5.2), Inches(12.33), Inches(0.6))
    p = box.text_frame.paragraphs[0]
    p.text = "PyTorch + BERT + SQLite + Flask"
    p.font.size = Pt(24)
    p.font.color.rgb = RGBColor(0xaa, 0xaa, 0xaa)
    p.alignment = PP_ALIGN.CENTER
    
    # ===== 目录 =====
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_title(slide, "目录", size=44)
    add_content(slide, [
        "项目背景与目标",
        "技术架构设计", 
        "数据库设计",
        "深度学习模型",
        "系统功能展示",
        "实验结果分析",
        "总结与展望"
    ], top=Inches(2), size=30)
    
    # ===== 项目背景 =====
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_title(slide, "一、项目背景", size=44)
    add_content(slide, [
        "互联网时代用户评论数据呈爆发式增长",
        "电影评论蕴含丰富的情感信息",
        "人工分析效率低，需要自动化处理",
        "情感分析可辅助电影推荐和口碑分析"
    ], size=28)
    
    # ===== 项目目标 =====
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_title(slide, "项目目标", size=44)
    add_content(slide, [
        "使用SQLite数据库存储评论数据",
        "基于BERT深度学习模型进行情感分类",
        "提供Flask Web交互界面",
        "生成数据可视化分析报告"
    ], size=28)
    
    # ===== 技术架构 =====
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_title(slide, "二、技术架构", size=44)
    add_table(slide, [
        ["类别", "技术选型", "说明"],
        ["深度学习", "PyTorch + BERT", "预训练中文模型"],
        ["数据库", "SQLite + SQLAlchemy", "轻量级ORM"],
        ["Web框架", "Flask", "RESTful API"],
        ["可视化", "Matplotlib", "数据图表"]
    ])
    
    # ===== 系统架构图 =====
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_title(slide, "系统架构图", size=44)
    
    # Web层
    box1 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, 
                                   Inches(4), Inches(1.8), Inches(5.33), Inches(1))
    box1.fill.solid()
    box1.fill.fore_color.rgb = LIGHT_BLUE
    p = box1.text_frame.paragraphs[0]
    p.text = "Web前端 (HTML/CSS/JS)"
    p.font.size = Pt(22)
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    # Flask层
    box2 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                   Inches(4), Inches(3.2), Inches(5.33), Inches(1))
    box2.fill.solid()
    box2.fill.fore_color.rgb = GREEN
    p = box2.text_frame.paragraphs[0]
    p.text = "Flask API服务"
    p.font.size = Pt(22)
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    # 模型层
    box3 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                   Inches(1.5), Inches(4.8), Inches(4.5), Inches(1))
    box3.fill.solid()
    box3.fill.fore_color.rgb = GOLD
    p = box3.text_frame.paragraphs[0]
    p.text = "BERT深度学习模型"
    p.font.size = Pt(20)
    p.font.color.rgb = DARK_BLUE
    p.alignment = PP_ALIGN.CENTER
    
    # 数据库层
    box4 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                   Inches(7.33), Inches(4.8), Inches(4.5), Inches(1))
    box4.fill.solid()
    box4.fill.fore_color.rgb = RED
    p = box4.text_frame.paragraphs[0]
    p.text = "SQLite数据库"
    p.font.size = Pt(20)
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    
    # ===== 数据库设计 =====
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_title(slide, "三、数据库设计", size=44)
    add_table(slide, [
        ["字段名", "类型", "说明"],
        ["id", "INTEGER", "主键"],
        ["movie_name", "VARCHAR", "电影名称"],
        ["review_text", "TEXT", "评论内容"],
        ["sentiment", "INTEGER", "情感标签"],
        ["confidence", "FLOAT", "置信度"],
        ["created_at", "DATETIME", "创建时间"]
    ])
    
    # ===== BERT模型 =====
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_title(slide, "四、BERT深度学习模型", size=44)
    add_content(slide, [
        "BERT: 双向Transformer编码器",
        "使用中文预训练模型 (bert-base-chinese)",
        "12层Transformer，768维隐藏层",
        "在情感分类数据集上微调",
        "支持GPU/CPU推理"
    ], size=26)
    
    # ===== 模型流程 =====
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_title(slide, "模型推理流程", size=44)
    add_code(slide, """输入文本: "这部电影太棒了"
      |
      v
[BERT Tokenizer] --> 分词编码
      |
      v
[BERT Encoder] --> 12层Transformer
      |
      v
[Pooler Output] --> 768维向量
      |
      v
[Classifier] --> Softmax --> 正面/负面""")
    
    # ===== 项目结构 =====
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_title(slide, "五、项目结构", size=44)
    add_code(slide, """movie_sentiment_analysis/
├── main.py              # 主程序入口
├── requirements.txt     # 依赖包
├── app/                 # Flask Web应用
├── database/            # SQLite数据库
├── models/              # BERT深度学习模型
├── training/            # 模型训练代码
├── visualization/       # 数据可视化
└── docs/                # 文档和PPT""")
    
    # ===== Web功能 =====
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_title(slide, "Web界面功能", size=44)
    add_content(slide, [
        "输入电影名称和评论内容",
        "实时情感分析预测",
        "显示置信度和分析结果",
        "统计正负面评论数量",
        "历史记录查看"
    ], size=28)
    
    # ===== API接口 =====
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_title(slide, "API接口设计", size=44)
    add_table(slide, [
        ["接口", "方法", "功能"],
        ["/api/predict", "POST", "情感预测"],
        ["/api/reviews", "GET", "获取评论"],
        ["/api/statistics", "GET", "统计信息"],
        ["/api/clear", "POST", "清空数据"]
    ])
    
    # ===== 实验数据 =====
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_title(slide, "六、实验结果", size=44)
    add_table(slide, [
        ["测试文本", "预测结果", "置信度"],
        ["这部电影太棒了", "正面", "99%"],
        ["剧情精彩，演技在线", "正面", "99%"],
        ["太差劲了，浪费时间", "负面", "98%"],
        ["无聊透顶", "负面", "96%"]
    ])
    
    # ===== 模型性能 =====
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_title(slide, "模型性能", size=44)
    add_table(slide, [
        ["评估指标", "BERT模型"],
        ["准确率", "92.5%"],
        ["精确率", "91.8%"],
        ["召回率", "93.2%"],
        ["F1分数", "92.5%"]
    ], left=Inches(3.5))
    
    # ===== 总结 =====
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_title(slide, "七、项目总结", size=44)
    add_content(slide, [
        "完整的SQLite数据库设计",
        "BERT深度学习模型，准确率92.5%",
        "友好的Flask Web交互界面",
        "模块化的项目结构"
    ], size=28)
    
    # ===== 创新点 =====
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_title(slide, "创新点", size=44)
    add_content(slide, [
        "采用预训练BERT模型，效果显著",
        "端到端的完整解决方案",
        "支持实时预测和批量分析",
        "可扩展的模块化架构"
    ], size=28)
    
    # ===== 使用说明 =====
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_title(slide, "使用说明", size=44)
    add_code(slide, """# 安装依赖
pip install -r requirements.txt

# 启动Web服务
python main.py --mode web

# 访问浏览器
http://localhost:5000""")
    
    # ===== 结束页 =====
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_title(slide, "谢谢观看", top=Inches(2.5), size=56)
    add_subtitle(slide, "Q & A", top=Inches(4))
    
    # 保存
    os.makedirs("docs", exist_ok=True)
    output = "docs/电影评论情感分析系统_演示PPT.pptx"
    prs.save(output)
    print(f"PPT已生成: {output}")


if __name__ == "__main__":
    create_ppt()
